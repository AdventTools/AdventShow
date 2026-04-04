#!/usr/bin/env python3
"""
Analyze PPT/PPTX files from the Exploratori collection.

Extracts raw text blocks from each slide of each file, detects patterns
(title structure, refrain markers, stanza numbers, master slide data),
and prints a summary of findings.

Uses:
  - python-pptx for .pptx files
  - olefile + binary CFB parsing for .ppt files (same approach as import.ts)
"""

import os
import re
import sys
import struct
import json
from pathlib import Path
from collections import Counter, defaultdict

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    import olefile
except ImportError:
    print("olefile not installed. Run: pip install olefile")
    sys.exit(1)


# ── Constants for PPT binary parsing (mirrors import.ts) ──────────────
PPT_RECORD_TYPE_SLIDE_PERSIST_ATOM = 1011
PPT_RECORD_TYPE_TEXT_HEADER_ATOM = 3999
PPT_RECORD_TYPE_TEXT_CHARS_ATOM = 4000
PPT_RECORD_TYPE_TEXT_BYTES_ATOM = 4008
PPT_RECORD_TYPE_SLIDE_LIST_WITH_TEXT = 4080

PPT_TEXT_TYPE_TITLE = 0
PPT_TEXT_TYPE_BODY = 1
PPT_TEXT_TYPE_CENTER_BODY = 5
PPT_TEXT_TYPE_CENTER_TITLE = 6
PPT_TEXT_TYPE_HALF_BODY = 7
PPT_TEXT_TYPE_QUARTER_BODY = 8

TEXT_TYPE_NAMES = {
    0: "Title",
    1: "Body",
    2: "Notes",
    3: "NotUsed",
    4: "Other",
    5: "CenterBody",
    6: "CenterTitle",
    7: "HalfBody",
    8: "QuarterBody",
}

MASTER_RE = re.compile(r"Click to edit Master")
STANZA_RE = re.compile(r"^(\d+)\.\s")
TITLE_NUM_RE = re.compile(r"^(\d+)\.\s+(.+)")
REFRAIN_LABEL_RE = re.compile(r"^(R\.|Refren:|Refrain:)", re.IGNORECASE)
REFRAIN_EXACT_RE = re.compile(r"^Refren:$", re.IGNORECASE)


# ── PPT binary parsing ───────────────────────────────────────────────
def read_record_header(data, offset):
    """Read an 8-byte PPT record header. Returns dict or None."""
    if offset + 8 > len(data):
        return None
    rec_word = struct.unpack_from('<H', data, offset)[0]
    rec_ver = rec_word & 0x000F
    rec_instance = rec_word >> 4
    rec_type = struct.unpack_from('<H', data, offset + 2)[0]
    rec_len = struct.unpack_from('<I', data, offset + 4)[0]
    content_start = offset + 8
    content_end = content_start + rec_len
    if content_end > len(data):
        return None
    return {
        'rec_ver': rec_ver,
        'rec_instance': rec_instance,
        'rec_type': rec_type,
        'rec_len': rec_len,
        'content_start': content_start,
        'content_end': content_end,
    }


def walk_records(data, start, end, visitor):
    """Walk PPT records, calling visitor(record) for each."""
    offset = start
    while offset + 8 <= end:
        rec = read_record_header(data, offset)
        if rec is None or rec['content_end'] > end:
            break
        visitor(rec)
        if rec['rec_ver'] == 0x0F:
            walk_records(data, rec['content_start'], rec['content_end'], visitor)
        offset = rec['content_end']


def find_slide_list_content(data):
    """Find the SlideListWithText container in the PPT binary."""
    result = [None]
    def visitor(rec):
        if (rec['rec_type'] == PPT_RECORD_TYPE_SLIDE_LIST_WITH_TEXT
                and rec['rec_instance'] == 0
                and rec['rec_ver'] == 0x0F):
            result[0] = data[rec['content_start']:rec['content_end']]
    walk_records(data, 0, len(data), visitor)
    return result[0]


def normalize_ppt_text(text):
    """Clean up text from PPT binary extraction."""
    text = text.replace('\x00', '')
    text = text.replace('\x0b', '\n')
    text = text.replace('\r\n', '\n')
    text = text.replace('\r', '\n')
    lines = [l.strip() for l in text.split('\n')]
    lines = [l for l in lines if l]
    return '\n'.join(lines).strip()


def extract_ppt_slides(filepath):
    """
    Extract slides from a .ppt file using CFB/OLE + binary record parsing.
    Returns list of slides, each slide is a list of (text_type, raw_text) tuples.
    """
    try:
        ole = olefile.OleFileIO(filepath)
    except Exception as e:
        return None, f"OLE open error: {e}"

    # Find the PowerPoint Document stream
    stream_name = None
    for candidate in ['PowerPoint Document', 'PP97_DUALSTORAGE/PowerPoint Document']:
        if ole.exists(candidate):
            stream_name = candidate
            break

    if not stream_name:
        ole.close()
        return None, "No 'PowerPoint Document' stream found"

    data = ole.openstream(stream_name).read()
    ole.close()

    slide_list = find_slide_list_content(data)
    if slide_list is None:
        return None, "No SlideListWithText found"

    slides = []
    current_slide = None
    current_text_type = None
    current_text_chunks = []

    def flush_text():
        nonlocal current_text_chunks
        if current_slide is not None and current_text_type is not None and current_text_chunks:
            raw = ''.join(current_text_chunks)
            text = normalize_ppt_text(raw)
            if text:
                current_slide.append((current_text_type, text))
        current_text_chunks = []

    offset = 0
    while offset + 8 <= len(slide_list):
        rec = read_record_header(slide_list, offset)
        if rec is None:
            break

        if rec['rec_type'] == PPT_RECORD_TYPE_SLIDE_PERSIST_ATOM:
            flush_text()
            current_text_type = None
            current_slide = []
            slides.append(current_slide)

        elif rec['rec_type'] == PPT_RECORD_TYPE_TEXT_HEADER_ATOM:
            flush_text()
            if current_slide is not None and rec['rec_len'] >= 4:
                current_text_type = struct.unpack_from('<I', slide_list, rec['content_start'])[0]
            else:
                current_text_type = None

        elif rec['rec_type'] == PPT_RECORD_TYPE_TEXT_CHARS_ATOM:
            if current_slide is not None and current_text_type is not None:
                current_text_chunks.append(
                    slide_list[rec['content_start']:rec['content_end']].decode('utf-16-le', errors='replace')
                )

        elif rec['rec_type'] == PPT_RECORD_TYPE_TEXT_BYTES_ATOM:
            if current_slide is not None and current_text_type is not None:
                current_text_chunks.append(
                    slide_list[rec['content_start']:rec['content_end']].decode('latin-1', errors='replace')
                )

        offset = rec['content_end']

    flush_text()

    return slides, None


def extract_pptx_slides(filepath):
    """
    Extract slides from a .pptx file using python-pptx.
    Returns list of slides, each slide is a list of (text_type_str, raw_text) tuples.
    """
    try:
        prs = Presentation(filepath)
    except Exception as e:
        return None, f"PPTX open error: {e}"

    slides = []
    for slide in prs.slides:
        slide_blocks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in para.runs)
                    if text.strip():
                        paragraphs.append(text.strip())
                if paragraphs:
                    combined = '\n'.join(paragraphs)
                    # Guess type from shape properties
                    shape_type = "Body"
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                        ph_idx = shape.placeholder_format.idx
                        if ph_idx == 0:
                            shape_type = "Title"
                        elif ph_idx == 1:
                            shape_type = "Body"
                    slide_blocks.append((shape_type, combined))
        slides.append(slide_blocks)

    return slides, None


# ── Pattern detection ─────────────────────────────────────────────────
def analyze_blocks(slides, filename):
    """Analyze a single file's slides for patterns."""
    info = {
        'filename': filename,
        'num_slides': len(slides),
        'has_master_blocks': False,
        'master_block_count': 0,
        'has_stanza_numbers': False,
        'stanza_numbers_found': [],
        'has_refrain_marker': False,
        'refrain_markers': [],
        'has_title_with_number': False,
        'title_pattern': None,
        'has_star_blocks': False,
        'has_footer_blocks': False,
        'total_text_blocks': 0,
        'text_types_seen': Counter(),
        'all_blocks': [],  # for detailed output
        'issues': [],
    }

    for si, slide in enumerate(slides):
        slide_info = {'slide_idx': si, 'blocks': []}
        for text_type, text in slide:
            type_name = TEXT_TYPE_NAMES.get(text_type, str(text_type)) if isinstance(text_type, int) else text_type
            info['text_types_seen'][type_name] += 1
            info['total_text_blocks'] += 1

            block_info = {
                'type': type_name,
                'text': text,
                'text_preview': text[:120].replace('\n', '\\n'),
            }
            slide_info['blocks'].append(block_info)

            # Check for master slide boilerplate
            if MASTER_RE.search(text):
                info['has_master_blocks'] = True
                info['master_block_count'] += 1

            # Check for star placeholder
            if text.strip() == '*':
                info['has_star_blocks'] = True

            # Check for footer-like text (e.g., "1/4")
            if re.match(r'^\d+/\d+$', text.strip()):
                info['has_footer_blocks'] = True

            # Check for stanza numbers
            m = STANZA_RE.match(text.strip())
            if m:
                info['has_stanza_numbers'] = True
                snum = int(m.group(1))
                if snum not in info['stanza_numbers_found']:
                    info['stanza_numbers_found'].append(snum)

            # Check for refrain markers
            if REFRAIN_EXACT_RE.match(text.strip()):
                info['has_refrain_marker'] = True
                info['refrain_markers'].append('Refren:')
            elif re.match(r'^R\.\s', text.strip()):
                info['has_refrain_marker'] = True
                info['refrain_markers'].append('R.')
            elif text.strip().startswith('Refren:') and len(text.strip()) > 10:
                info['has_refrain_marker'] = True
                info['refrain_markers'].append('Refren: (inline)')

            # Check for title with "Imnul NNN" or "NNN. Title" pattern
            flat_text = re.sub(r'\s+', ' ', text.replace('\n', ' ')).strip()
            if re.match(r'Imnul\s+\d+', flat_text, re.IGNORECASE):
                info['has_title_with_number'] = True
                info['title_pattern'] = f'"Imnul NNN": {flat_text[:80]}'
            elif TITLE_NUM_RE.match(flat_text) and len(flat_text) < 80:
                if not info['has_title_with_number']:
                    info['has_title_with_number'] = True
                    info['title_pattern'] = f'"NNN. Title": {flat_text[:80]}'

        info['all_blocks'].append(slide_info)

    return info


# ── Pretty printing ──────────────────────────────────────────────────
def print_file_detail(info, verbose=True):
    """Print detailed analysis of one file."""
    print(f"\n{'='*72}")
    print(f"FILE: {info['filename']}")
    print(f"  Slides: {info['num_slides']}")
    print(f"  Total text blocks: {info['total_text_blocks']}")
    print(f"  Text types: {dict(info['text_types_seen'])}")
    print(f"  Master slide blocks: {info['has_master_blocks']} (count: {info['master_block_count']})")
    print(f"  Star (*) blocks: {info['has_star_blocks']}")
    print(f"  Footer blocks: {info['has_footer_blocks']}")
    print(f"  Title with number: {info['has_title_with_number']} → {info['title_pattern']}")
    print(f"  Stanza numbers: {info['has_stanza_numbers']} → {info['stanza_numbers_found']}")
    print(f"  Refrain markers: {info['has_refrain_marker']} → {info['refrain_markers']}")

    if info['issues']:
        for issue in info['issues']:
            print(f"  ⚠️  {issue}")

    if verbose:
        for slide in info['all_blocks']:
            print(f"\n  --- Slide {slide['slide_idx']} ---")
            for bi, block in enumerate(slide['blocks']):
                marker = ""
                if MASTER_RE.search(block['text']):
                    marker = " [MASTER]"
                elif block['text'].strip() == '*':
                    marker = " [STAR]"
                print(f"    Block {bi} ({block['type']}){marker}: {block['text_preview']}")


def print_summary(all_info):
    """Print aggregate summary across all files."""
    print(f"\n\n{'#'*72}")
    print(f"# AGGREGATE SUMMARY across {len(all_info)} files")
    print(f"{'#'*72}")

    total = len(all_info)
    with_master = sum(1 for i in all_info if i['has_master_blocks'])
    with_star = sum(1 for i in all_info if i['has_star_blocks'])
    with_footer = sum(1 for i in all_info if i['has_footer_blocks'])
    with_title_num = sum(1 for i in all_info if i['has_title_with_number'])
    with_stanza_nums = sum(1 for i in all_info if i['has_stanza_numbers'])
    with_refrain = sum(1 for i in all_info if i['has_refrain_marker'])

    slide_counts = Counter(i['num_slides'] for i in all_info)
    all_types = Counter()
    for i in all_info:
        all_types.update(i['text_types_seen'])

    all_refrain_types = Counter()
    for i in all_info:
        for rm in i['refrain_markers']:
            all_refrain_types[rm] += 1

    print(f"\n  Files with master blocks ('Click to edit Master'): {with_master}/{total}")
    print(f"  Files with star (*) blocks: {with_star}/{total}")
    print(f"  Files with footer blocks (N/M): {with_footer}/{total}")
    print(f"  Files with title+number pattern: {with_title_num}/{total}")
    print(f"  Files with stanza numbers: {with_stanza_nums}/{total}")
    print(f"  Files with refrain markers: {with_refrain}/{total}")

    print(f"\n  Slide count distribution:")
    for sc in sorted(slide_counts.keys()):
        print(f"    {sc} slides: {slide_counts[sc]} files")

    print(f"\n  Text type distribution (across all blocks):")
    for tt, count in all_types.most_common():
        print(f"    {tt}: {count}")

    print(f"\n  Refrain marker types:")
    for rm, count in all_refrain_types.most_common():
        print(f"    {rm}: {count}")

    # Edge cases
    no_stanza = [i['filename'] for i in all_info if not i['has_stanza_numbers'] and i['num_slides'] <= 2]
    many_slides = [i['filename'] for i in all_info if i['num_slides'] >= 8]
    no_title = [i['filename'] for i in all_info if not i['has_title_with_number']]

    if no_stanza:
        print(f"\n  Files without stanza numbers (≤2 slides): {len(no_stanza)}")
        for f in no_stanza[:10]:
            print(f"    - {f}")
        if len(no_stanza) > 10:
            print(f"    ... and {len(no_stanza) - 10} more")

    if many_slides:
        print(f"\n  Files with many slides (≥8): {len(many_slides)}")
        for f in many_slides:
            print(f"    - {f}")

    if no_title:
        print(f"\n  Files without detected title+number: {len(no_title)}")
        for f in no_title[:10]:
            print(f"    - {f}")
        if len(no_title) > 10:
            print(f"    ... and {len(no_title) - 10} more")


# ── Main ──────────────────────────────────────────────────────────────
def main():
    import argparse
    parser = argparse.ArgumentParser(description='Analyze Exploratori PPT/PPTX files')
    parser.add_argument('directory', help='Directory containing PPT/PPTX files')
    parser.add_argument('--limit', type=int, default=0, help='Limit number of files to analyze (0=all)')
    parser.add_argument('--verbose', action='store_true', help='Show detailed block output')
    parser.add_argument('--summary-only', action='store_true', help='Only show summary')
    parser.add_argument('--json-out', type=str, default='', help='Output raw data as JSON to file')
    args = parser.parse_args()

    directory = Path(args.directory)
    if not directory.is_dir():
        print(f"Error: {directory} is not a directory")
        sys.exit(1)

    # Collect files, sorted by name
    files = sorted([
        f for f in directory.iterdir()
        if f.suffix.lower() in ('.ppt', '.pptx')
    ], key=lambda f: f.name)

    if args.limit > 0:
        files = files[:args.limit]

    print(f"Found {len(files)} PPT/PPTX files in {directory}")
    print(f"Analyzing {'first ' + str(args.limit) if args.limit else 'all'} files...\n")

    all_info = []
    raw_data = {}  # for JSON export

    for filepath in files:
        fname = filepath.name
        ext = filepath.suffix.lower()

        if ext == '.ppt':
            slides, error = extract_ppt_slides(str(filepath))
        elif ext == '.pptx':
            slides, error = extract_pptx_slides(str(filepath))
        else:
            continue

        if error:
            print(f"  ERROR processing {fname}: {error}")
            continue

        if slides is None:
            print(f"  ERROR: No slides extracted from {fname}")
            continue

        info = analyze_blocks(slides, fname)
        all_info.append(info)

        # Store raw data
        raw_slides = []
        for slide in slides:
            raw_slides.append([
                {'type': TEXT_TYPE_NAMES.get(tt, str(tt)) if isinstance(tt, int) else tt, 'text': text}
                for tt, text in slide
            ])
        raw_data[fname] = {
            'num_slides': len(slides),
            'slides': raw_slides,
        }

        if not args.summary_only:
            print_file_detail(info, verbose=args.verbose)

    if all_info:
        print_summary(all_info)

    if args.json_out:
        with open(args.json_out, 'w', encoding='utf-8') as f:
            json.dump(raw_data, f, ensure_ascii=False, indent=2)
        print(f"\nRaw data written to {args.json_out}")


if __name__ == '__main__':
    main()
